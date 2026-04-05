import logging
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from .chart_mapper import ChartMapper

logger = logging.getLogger(__name__)

class PPTXExporter:
    """
    The 'Rigid Layout Engine' for Consulting-Grade Deliverables.
    Enforces a strict visual grid and premium branding (#8B5CF6).
    Structure: Title -> Overview -> Summary -> Insights -> Appendix.
    """
    def __init__(self):
        self.prs = Presentation()
        # Slide Layouts: 0=Title, 1=Title+Content, 5=Blank (Custom)
        self.ACCENT_COLOR = RGBColor(0x8B, 0x5C, 0xF6) # Violet 500

    def _set_slide_title(self, slide, title_text):
        """Standardized title placement top-left."""
        title = slide.shapes.title
        title.text = title_text
        title.left = Inches(0.5)
        title.top = Inches(0.2)
        title.width = Inches(9.0)
        title.height = Inches(0.6)
        
        # Style
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.ACCENT_COLOR
        p.alignment = PP_ALIGN.LEFT

    def create_title_slide(self, context: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = context['session_title']
        subtitle = slide.placeholders[1]
        subtitle.text = f"Consulting Analysis Report\nGenerated on {context['date']}\nData: {context['dataset_info']['name']}"

    def create_overview_slide(self, context: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_slide_title(slide, "Dataset Overview")
        content = slide.placeholders[1]
        info = context['dataset_info']
        text = f"• Total Rows: {info['rows']:,}\n"
        text += f"• Total Columns: {info['cols']}\n"
        text += f"• Key Fields: {', '.join(info['fields'])}\n"
        content.text = text

    def create_summary_slide(self, context: dict):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        self._set_slide_title(slide, "Executive Summary")
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in context['executive_summary']:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(10)

    def create_insight_slide(self, insight: dict):
        """
        Follows the Narrative Grid: [Title] -> [Chart] -> [3 Takeaways].
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5]) # Blank
        self._set_slide_title(slide, insight['title'])
        
        # 1. Visualization (Chart/Image)
        # We try native chart first, then image fallback
        visual_placed = False
        for artifact in insight.get('artifacts', []):
            if artifact.get('type') == 'chart':
                mapping = ChartMapper.map_to_pptx_chart(artifact)
                if mapping:
                    chart_type, chart_data = mapping
                    x, y, cx, cy = Inches(0.5), Inches(1.0), Inches(5.5), Inches(4.5)
                    slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
                    visual_placed = True
                    break
                else:
                    # Fallback Image
                    img_stream = ChartMapper.generate_static_image(artifact)
                    if img_stream:
                        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.0), width=Inches(5.5))
                        visual_placed = True
                        break

        # 2. Key Takeaways (Bullets on the Right)
        tx_box = slide.shapes.add_textbox(Inches(6.2), Inches(1.0), Inches(3.3), Inches(4.5))
        tf = tx_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "Key Takeaways:"
        p.font.bold = True
        p.font.size = Pt(16)
        
        for takeaway in insight.get('takeaways', []):
            p = tf.add_paragraph()
            p.text = takeaway
            p.font.size = Pt(14)
            p.level = 0
            p.space_before = Pt(8)

    def generate_report(self, context: dict) -> io.BytesIO:
        """Orchestrates slide generation."""
        self.create_title_slide(context)
        self.create_overview_slide(context)
        self.create_summary_slide(context)
        
        for insight in context['insights']:
            self.create_insight_slide(insight)
            
        # Re-save to stream
        pptx_stream = io.BytesIO()
        self.prs.save(pptx_stream)
        pptx_stream.seek(0)
        return pptx_stream
